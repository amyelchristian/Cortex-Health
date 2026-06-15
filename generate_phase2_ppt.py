"""
generate_phase2_ppt.py
Generates a Phase 2 PowerPoint presentation for CORTEX v3.0
Matching the exact style/theme of the Phase 1 PPT.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
BG_PATH = os.path.join(BASE_DIR, "Assets", "bg.jpg")
HEADER_PATH = os.path.join(BASE_DIR, "PPT", "extracted_logos", "full_header.png")
OUTPUT_PATH = os.path.join(BASE_DIR, "PPT", "Phase 2 ppt.pptx")

# Colors matching Phase 1 theme
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
MUTED_GREEN = RGBColor(150, 200, 150)
DARK_BG = RGBColor(20, 30, 20)
CARD_BG = RGBColor(40, 50, 40)
ACCENT_GREEN = RGBColor(0, 218, 170)
LABEL_BG = RGBColor(80, 60, 60)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, prs):
    """Add bg.jpg as full-slide background image."""
    slide.shapes.add_picture(BG_PATH, Emu(0), Emu(0), prs.slide_width, prs.slide_height)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    """Helper to add a text box with standard formatting."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox, tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=0):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p


def add_card(slide, left, top, width, height, opacity=40):
    """Add a semi-transparent rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(40, 55, 40)
    shape.fill.fore_color.brightness = 0.1
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_label(slide, left, top, text, width=2.0):
    """Add a small label/tag chip like in Phase 1."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(0.35)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(60, 50, 50)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return shape


# ─────────────────────────── SLIDES ───────────────────────────

def slide_cover(prs):
    """Slide 1: Cover - identical layout to Phase 1 but Phase 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    add_bg(slide, prs)

    # Header strip with logos (PDEU, Gemini, GSA)
    slide.shapes.add_picture(HEADER_PATH, Emu(0), Emu(0),
                             prs.slide_width, Inches(1.05))

    # "Cover Slide" label
    add_label(slide, 5.5, 1.3, "Cover Slide", width=1.5)

    # Title
    add_text_box(slide, 1.5, 1.8, 10, 1.5,
                 "Predicting Patient Health Risk\nThrough Proactive Monitoring",
                 font_size=44, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, 2, 3.4, 9, 0.6,
                 "An Early Warning Approach to Prevent Patient Deterioration",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name="Courier New")

    # Team info block
    info_lines = [
        ("College:", "Pandit Deendayal Energy University"),
        ("Team Name:", "Tech Titans"),
        ("Team Members:", "Amyel Christian, Malhar Gajjar, Darsh Patel"),
        ("Event:", "India Tech Summit - Innovate 2026"),
        ("Phase:", "Ideathon - Phase 2"),
        ("Team ID:", "8HmuRy2kMbeU8JQJhy0E"),
        ("Problem Statement ID:", "Xfair1003"),
    ]
    y = 4.6
    for label, value in info_lines:
        txBox, tf = add_text_box(slide, 0.5, y, 10, 0.35, "", font_size=14)
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = label + " "
        run1.font.bold = True
        run1.font.size = Pt(14)
        run1.font.color.rgb = WHITE
        run1.font.name = "Calibri"
        run2 = p.add_run()
        run2.text = value
        run2.font.bold = False
        run2.font.size = Pt(14)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"
        y += 0.35


def slide_section_title(prs, title, subtitle):
    """Section title slide - large centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_text_box(slide, 1, 2.0, 11, 2.5, title,
                 font_size=56, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 4.2, 9, 1.0, subtitle,
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_problem_recap(prs):
    """Slide 3: Problem recap from Phase 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.4, 6, 0.8, "PROBLEM OVERVIEW",
                 font_size=40, bold=True)

    add_text_box(slide, 0.5, 1.4, 6.5, 2.5,
                 "Despite the availability of rich patient data, current "
                 "healthcare systems lack intelligent, real-time analytics "
                 "to predict patient deterioration.\n\n"
                 "The absence of early warning mechanisms limits proactive care "
                 "and reduces the ability to intervene before conditions become "
                 "life-threatening.",
                 font_size=16, color=LIGHT_GRAY)

    # Right side cards
    card_data = [
        ("Unutilized\nPatient Data", 8.0, 0.8),
        ("Reactive Healthcare\nApproach", 8.0, 2.8),
        ("Proactive Risk\nAssessment Opportunity", 8.0, 4.8),
    ]
    for text, x, y in card_data:
        card = add_card(slide, x, y, 4.5, 1.5)
        add_text_box(slide, x + 0.3, y + 0.3, 3.8, 1.0, text,
                     font_size=18, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 0.5, 5.8, 7, 1.0,
                 "Healthcare lacks intelligent systems that convert existing "
                 "patient data into timely, actionable insights, resulting in "
                 "delayed interventions and preventable critical events.",
                 font_size=14, color=LIGHT_GRAY, bold=True)


def slide_solution_overview(prs):
    """Slide 4: Our Solution - CORTEX v3.0."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "Our Solution: CORTEX v3.0",
                 font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.1, 12, 0.5,
                 "AI-Powered Clinical Risk Assessment System - Fully Deployed & Live",
                 font_size=16, color=MUTED_GREEN)

    # 4 feature cards in 2x2 grid
    features = [
        ("ML Risk Prediction", "XGBoost model trained on 1.41M+ MIMIC-IV\n"
         "patient records with 99.98% accuracy.\n"
         "Real-time risk scoring: Low / Medium / High."),
        ("AI Health Companion", "LLM-powered health chat using Llama 3.3 70B\n"
         "via OpenRouter API. Context-aware\n"
         "conversations with assessment data."),
        ("Document Intelligence", "Upload and analyze clinical documents\n"
         "(lab reports, imaging, prescriptions).\n"
         "AI-powered OCR and summarization."),
        ("Safety Override System", "Rule-based critical threshold detection.\n"
         "SpO2 < 90%, HR > 150, BP > 180/110\n"
         "triggers immediate high-risk alerts."),
    ]
    positions = [(0.5, 2.0), (6.8, 2.0), (0.5, 4.8), (6.8, 4.8)]
    for (title, desc), (x, y) in zip(features, positions):
        add_card(slide, x, y, 5.8, 2.3)
        add_text_box(slide, x + 0.3, y + 0.25, 5.2, 0.5, title,
                     font_size=20, bold=True, color=WHITE)
        add_text_box(slide, x + 0.3, y + 0.9, 5.2, 1.3, desc,
                     font_size=13, color=LIGHT_GRAY)


def slide_tech_stack(prs):
    """Slide 5: Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "Technology Stack", font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.1, 12, 0.5,
                 "Production-grade technologies powering CORTEX",
                 font_size=16, color=LIGHT_GRAY)

    categories = [
        ("Frontend", [
            "React 19.2.0 - Component-based UI",
            "Vite 7.3.1 - Next-gen build tool",
            "Tailwind CSS 4.2 - Utility-first styling",
            "React Router 7.13 - Client-side routing",
            "Lucide React - Icon library",
        ]),
        ("Backend", [
            "FastAPI - High-performance Python API",
            "Uvicorn - ASGI server",
            "Python 3.11 - Runtime environment",
            "Docker - Containerization",
            "Google Cloud Run - Serverless hosting",
        ]),
        ("AI / ML", [
            "XGBoost - Primary ML model",
            "scikit-learn - ML pipeline & evaluation",
            "MIMIC-IV - Clinical training data",
            "OpenRouter API - LLM gateway",
            "Llama 3.3 70B - Health chat model",
        ]),
        ("Infrastructure", [
            "Firebase Auth - User authentication",
            "Cloud Firestore - NoSQL database",
            "Vercel - Frontend hosting",
            "Google Cloud Run - API hosting",
            "Firebase Security Rules",
        ]),
    ]

    x_positions = [0.3, 3.4, 6.5, 9.6]
    for i, (cat_title, items) in enumerate(categories):
        x = x_positions[i]
        add_card(slide, x, 1.8, 2.9, 5.2)
        add_text_box(slide, x + 0.15, 1.95, 2.6, 0.5, cat_title,
                     font_size=18, bold=True, color=ACCENT_GREEN,
                     alignment=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text_box(slide, x + 0.15, 2.6 + j * 0.7, 2.6, 0.65, item,
                         font_size=11, color=LIGHT_GRAY)


def slide_system_architecture(prs):
    """Slide 6: System Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "System Architecture", font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.0, 5, 0.5,
                 "TECHNICAL ARCHITECTURE",
                 font_size=12, color=LIGHT_GRAY, font_name="Courier New")

    # Architecture flow cards
    layers = [
        ("User Interface Layer", "React 19 SPA\nDark glass-morphism design\n"
         "6 dashboard tabs\nResponsive & mobile-first", 0.5, 1.8),
        ("API Gateway", "FastAPI on Cloud Run\nRESTful endpoints\n"
         "CORS middleware\nMultipart file handling", 4.5, 1.8),
        ("Intelligence Layer", "CortexUnifiedAgent\nRoutes to ML / LLM / DB\n"
         "Safety override checks\nContext management", 8.5, 1.8),
    ]
    for title, desc, x, y in layers:
        add_card(slide, x, y, 3.5, 2.2)
        add_text_box(slide, x + 0.2, y + 0.15, 3.1, 0.5, title,
                     font_size=16, bold=True)
        add_text_box(slide, x + 0.2, y + 0.7, 3.1, 1.4, desc,
                     font_size=12, color=LIGHT_GRAY)

    # Bottom row
    bottom_layers = [
        ("ML Engine", "XGBoost model\n44 engineered features\n"
         "Risk classification\nFeature importance", 0.5, 4.5),
        ("LLM Service", "Llama 3.3 70B\nHealth chat responses\n"
         "Assessment explanations\nDocument analysis", 4.5, 4.5),
        ("Data Layer", "Cloud Firestore\nFirebase Auth\n"
         "Per-user collections\nReal-time sync", 8.5, 4.5),
    ]
    for title, desc, x, y in bottom_layers:
        add_card(slide, x, y, 3.5, 2.2)
        add_text_box(slide, x + 0.2, y + 0.15, 3.1, 0.5, title,
                     font_size=16, bold=True)
        add_text_box(slide, x + 0.2, y + 0.7, 3.1, 1.4, desc,
                     font_size=12, color=LIGHT_GRAY)

    # Arrows between top cards
    for x in [4.0, 8.0]:
        add_text_box(slide, x, 2.5, 0.5, 0.5, "->",
                     font_size=24, bold=True, color=ACCENT_GREEN,
                     alignment=PP_ALIGN.CENTER)


def slide_ml_model(prs):
    """Slide 7: ML Model Details."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "ML Model: Clinical Risk Prediction",
                 font_size=40, bold=True)

    # Left side - model stats
    add_card(slide, 0.5, 1.4, 5.8, 5.5)
    stats = [
        ("Algorithm", "XGBoost (Gradient Boosted Trees)"),
        ("Training Data", "MIMIC-IV Clinical Database"),
        ("Training Samples", "1,410,937 patient records"),
        ("Features", "44 engineered clinical features"),
        ("Accuracy", "99.98%"),
        ("High-Risk Recall", "100% (never misses critical patients)"),
        ("Inference Time", "< 100ms per prediction"),
        ("Risk Categories", "Low / Medium / High"),
        ("Threshold", "0.15 (optimized for sensitivity)"),
        ("Safety Overrides", "Rule-based critical thresholds"),
    ]
    for i, (label, value) in enumerate(stats):
        y = 1.6 + i * 0.5
        txBox, tf = add_text_box(slide, 0.8, y, 5.2, 0.45, "", font_size=13)
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = label + ":  "
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = ACCENT_GREEN
        run1.font.name = "Calibri"
        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(13)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"

    # Right side - key features
    add_card(slide, 6.8, 1.4, 5.8, 5.5)
    add_text_box(slide, 7.1, 1.5, 5.2, 0.5,
                 "Core Vital Signs Monitored",
                 font_size=18, bold=True, color=ACCENT_GREEN)
    vitals = [
        "Heart Rate (HR) - bpm",
        "Systolic Blood Pressure (SBP) - mmHg",
        "Diastolic Blood Pressure (DBP) - mmHg",
        "Oxygen Saturation (SpO2) - %",
        "Respiratory Rate (RR) - breaths/min",
        "Body Temperature - Fahrenheit",
    ]
    for i, v in enumerate(vitals):
        add_text_box(slide, 7.3, 2.2 + i * 0.45, 5.0, 0.4,
                     "  " + v, font_size=13, color=LIGHT_GRAY)

    add_text_box(slide, 7.1, 5.0, 5.2, 0.5,
                 "Derived Features",
                 font_size=18, bold=True, color=ACCENT_GREEN)
    derived = [
        "Mean Arterial Pressure (MAP)",
        "Shock Index (HR / SBP)",
        "Pulse Pressure (SBP - DBP)",
        "Rolling averages & rate of change",
    ]
    for i, d in enumerate(derived):
        add_text_box(slide, 7.3, 5.5 + i * 0.4, 5.0, 0.35,
                     "  " + d, font_size=12, color=LIGHT_GRAY)


def slide_safety_overrides(prs):
    """Slide 8: Safety Override System."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "Safety Override System",
                 font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.1, 12, 0.5,
                 "Critical threshold detection that overrides ML predictions to ensure patient safety",
                 font_size=16, color=LIGHT_GRAY)

    overrides = [
        ("SpO2 < 90%", "Critical Oxygen\nDrop", "Immediate high-risk\nalert triggered"),
        ("HR > 150 or < 40", "Critical Heart\nRate", "Dangerous tachycardia\nor bradycardia"),
        ("SBP > 180 / DBP > 110", "Hypertensive\nCrisis", "Emergency blood\npressure levels"),
        ("SBP < 90 / DBP < 60", "Hypotensive\nCrisis", "Dangerously low\nblood pressure"),
        ("Temp > 103F / < 95F", "Critical\nTemperature", "Hyperthermia or\nhypothermia risk"),
    ]

    x_positions = [0.3, 2.8, 5.3, 7.8, 10.3]
    for i, (threshold, title, desc) in enumerate(overrides):
        x = x_positions[i]
        add_card(slide, x, 2.0, 2.3, 4.5)
        add_text_box(slide, x + 0.1, 2.2, 2.1, 0.6, threshold,
                     font_size=14, bold=True, color=RGBColor(239, 68, 68),
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 3.0, 2.1, 0.8, title,
                     font_size=16, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 4.2, 2.1, 0.8, desc,
                     font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_dashboard_features(prs):
    """Slide 9: Dashboard Features."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "Dashboard Features", font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.0, 12, 0.5,
                 "6 comprehensive tabs for complete health management",
                 font_size=16, color=LIGHT_GRAY)

    tabs = [
        ("Home Tab", "Patient overview with\nCortex AI score,\n"
         "health summary,\ncare timeline, and\nquick actions."),
        ("Health Chat", "AI-powered health\nassistant using\n"
         "Llama 3.3 70B.\nContext-aware\nconversations."),
        ("Assessment", "Real-time vital sign\ninput with instant\n"
         "ML risk predictions\nand explanations."),
        ("Documents", "Upload & analyze\nclinical documents.\n"
         "AI-powered OCR\nand summarization."),
        ("Database", "Assessment history\nwith biometric logs.\n"
         "Track trends over\ntime with timestamps."),
        ("Analytics", "Brain mapping,\nrisk gauges, vitals\n"
         "grid, and clinical\naction plans."),
    ]
    for i, (title, desc) in enumerate(tabs):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 1.8 + row * 2.8
        add_card(slide, x, y, 3.8, 2.4)
        add_text_box(slide, x + 0.2, y + 0.15, 3.4, 0.5, title,
                     font_size=18, bold=True, color=ACCENT_GREEN)
        add_text_box(slide, x + 0.2, y + 0.7, 3.4, 1.6, desc,
                     font_size=13, color=LIGHT_GRAY)


def slide_auth_security(prs):
    """Slide 10: Authentication & Security."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "Authentication & Security", font_size=42, bold=True)

    # Left card
    add_card(slide, 0.5, 1.5, 5.8, 5.5)
    add_text_box(slide, 0.8, 1.7, 5.2, 0.5,
                 "Firebase Authentication", font_size=20, bold=True,
                 color=ACCENT_GREEN)
    auth_items = [
        "Email/password sign-up and login",
        "Dedicated login & signup pages (/login, /signup)",
        "Password reset functionality",
        "Protected dashboard routes",
        "AuthContext with React Context API",
        "LocalStorage caching for instant auth restore",
        "Session persistence across page reloads",
    ]
    for i, item in enumerate(auth_items):
        add_text_box(slide, 1.0, 2.4 + i * 0.55, 5.0, 0.5,
                     "  " + item, font_size=13, color=LIGHT_GRAY)

    # Right card
    add_card(slide, 6.8, 1.5, 5.8, 5.5)
    add_text_box(slide, 7.1, 1.7, 5.2, 0.5,
                 "Cloud Firestore Security", font_size=20, bold=True,
                 color=ACCENT_GREEN)
    security_items = [
        "Per-user document isolation (users/{userId})",
        "Firestore security rules enforce ownership",
        "Read/write restricted to authenticated owner",
        "OTP verifications locked to Cloud Functions",
        "CORS middleware on FastAPI backend",
        "Environment variable secrets management",
        "HTTPS-only API communication",
    ]
    for i, item in enumerate(security_items):
        add_text_box(slide, 7.3, 2.4 + i * 0.55, 5.0, 0.5,
                     "  " + item, font_size=13, color=LIGHT_GRAY)


def slide_api_endpoints(prs):
    """Slide 11: API Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "API Architecture", font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.0, 12, 0.5,
                 "RESTful API built with FastAPI, deployed on Google Cloud Run",
                 font_size=16, color=LIGHT_GRAY)

    endpoints = [
        ("POST", "/chat", "Conversational health queries with context"),
        ("POST", "/assessment", "ML risk prediction from vital signs"),
        ("POST", "/upload-document", "Clinical document analysis via AI"),
        ("GET", "/documents/{uid}", "Retrieve user's document records"),
        ("GET", "/dashboard/{uid}", "Full dashboard data aggregation"),
        ("POST", "/reminder", "Create health/medication reminders"),
        ("GET", "/history/{uid}", "Conversation history retrieval"),
    ]

    add_card(slide, 0.5, 1.6, 12, 5.2)
    # Table header
    add_text_box(slide, 0.8, 1.8, 1.5, 0.4, "Method",
                 font_size=14, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, 2.8, 1.8, 3.5, 0.4, "Endpoint",
                 font_size=14, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, 6.8, 1.8, 5.5, 0.4, "Description",
                 font_size=14, bold=True, color=ACCENT_GREEN)

    for i, (method, endpoint, desc) in enumerate(endpoints):
        y = 2.4 + i * 0.6
        method_color = RGBColor(34, 197, 94) if method == "GET" else RGBColor(59, 130, 246)
        add_text_box(slide, 0.8, y, 1.5, 0.4, method,
                     font_size=13, bold=True, color=method_color,
                     font_name="Courier New")
        add_text_box(slide, 2.8, y, 3.5, 0.4, endpoint,
                     font_size=13, color=WHITE, font_name="Courier New")
        add_text_box(slide, 6.8, y, 5.5, 0.4, desc,
                     font_size=13, color=LIGHT_GRAY)


def slide_deployment(prs):
    """Slide 12: Deployment Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "Deployment Architecture", font_size=42, bold=True)

    # Three deployment cards
    deployments = [
        ("Frontend - Vercel", [
            "React SPA built with Vite 7",
            "Deployed on Vercel platform",
            "SPA rewrites for /login, /signup",
            "Asset caching (31536000s immutable)",
            "Environment: .env.production",
            "Auto-deploy on git push",
        ]),
        ("Backend - Google Cloud Run", [
            "Docker container (Python 3.11-slim)",
            "FastAPI with Uvicorn server",
            "Auto-scaling (0 to N instances)",
            "Port 8080 configuration",
            "Environment secrets via GCP",
            "Region: us-central1",
        ]),
        ("Database - Firebase", [
            "Cloud Firestore (NoSQL)",
            "Firebase Authentication",
            "Real-time data sync",
            "Security rules enforcement",
            "Per-user data isolation",
            "Automatic backups",
        ]),
    ]

    for i, (title, items) in enumerate(deployments):
        x = 0.3 + i * 4.3
        add_card(slide, x, 1.5, 3.9, 5.5)
        add_text_box(slide, x + 0.2, 1.7, 3.5, 0.5, title,
                     font_size=18, bold=True, color=ACCENT_GREEN,
                     alignment=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text_box(slide, x + 0.25, 2.5 + j * 0.6, 3.4, 0.55,
                         "  " + item, font_size=12, color=LIGHT_GRAY)


def slide_data_flow(prs):
    """Slide 13: Data Flow - Risk Assessment Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "Risk Assessment Pipeline", font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.0, 12, 0.5,
                 "End-to-end flow from vital sign input to clinical recommendation",
                 font_size=16, color=LIGHT_GRAY)

    steps = [
        ("1. Vital Input", "User enters 6 core\nvitals via dashboard\nAssessment Tab"),
        ("2. Safety Check", "Critical threshold\noverrides are checked\nbefore ML inference"),
        ("3. Feature Eng.", "6 vitals expanded to\n44 clinical features\n(MAP, Shock Index...)"),
        ("4. ML Prediction", "XGBoost classifies\nrisk as Low/Medium/\nHigh with confidence"),
        ("5. LLM Explain", "Llama 3.3 generates\nplain-language clinical\nexplanation"),
        ("6. Persist & Show", "Results saved to\nFirestore and displayed\non dashboard"),
    ]

    for i, (title, desc) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 1.8 + row * 2.8
        add_card(slide, x, y, 3.8, 2.3)
        add_text_box(slide, x + 0.2, y + 0.15, 3.4, 0.5, title,
                     font_size=18, bold=True, color=ACCENT_GREEN)
        add_text_box(slide, x + 0.2, y + 0.7, 3.4, 1.4, desc,
                     font_size=14, color=LIGHT_GRAY)


def slide_model_performance(prs):
    """Slide 14: Model Performance Metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "Model Performance", font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.0, 12, 0.5,
                 "Production model evaluation results on held-out test data",
                 font_size=16, color=LIGHT_GRAY)

    # Performance metric cards
    metrics = [
        ("99.98%", "Overall\nAccuracy"),
        ("100%", "High-Risk\nRecall"),
        ("< 100ms", "Inference\nLatency"),
        ("1.41M+", "Training\nSamples"),
    ]
    for i, (value, label) in enumerate(metrics):
        x = 0.5 + i * 3.2
        add_card(slide, x, 1.8, 2.8, 2.0)
        add_text_box(slide, x + 0.1, 1.95, 2.6, 0.8, value,
                     font_size=36, bold=True, color=ACCENT_GREEN,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 2.8, 2.6, 0.7, label,
                     font_size=16, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Model comparison table
    add_card(slide, 0.5, 4.2, 12, 3.0)
    add_text_box(slide, 0.8, 4.3, 5, 0.5,
                 "Model Comparison Results",
                 font_size=18, bold=True, color=ACCENT_GREEN)

    models = [
        ("Model", "Accuracy", "Precision", "Recall", "F1 Score"),
        ("XGBoost", "99.98%", "99.97%", "100%", "99.98%"),
        ("Random Forest", "99.95%", "99.93%", "99.96%", "99.94%"),
        ("Logistic Regression", "89.42%", "88.15%", "89.42%", "88.71%"),
    ]
    col_x = [0.8, 3.5, 5.5, 7.5, 9.5]
    for row_i, row in enumerate(models):
        y = 5.0 + row_i * 0.45
        is_header = row_i == 0
        for col_i, cell in enumerate(row):
            add_text_box(slide, col_x[col_i], y, 2.0, 0.4, cell,
                         font_size=13,
                         bold=is_header,
                         color=ACCENT_GREEN if is_header else LIGHT_GRAY,
                         font_name="Courier New" if not is_header else "Calibri")


def slide_ui_design(prs):
    """Slide 15: UI/UX Design Philosophy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 12, 0.8,
                 "UI/UX Design", font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.0, 12, 0.5,
                 "Premium dark glass-morphism design with clinical precision",
                 font_size=16, color=LIGHT_GRAY)

    design_elements = [
        ("Dark Glass Morphism", "Frosted glass cards with\nbackdrop blur effects.\n"
         "Semi-transparent overlays\ncreate depth and hierarchy."),
        ("Ambient Mesh Backgrounds", "Dynamic gradient mesh\nbackgrounds that adapt\n"
         "to risk levels. Green for\nlow, amber for medium,\nred for high risk."),
        ("Clinical Color System", "Risk-aware color coding:\n"
         "Green (#22C55E) - Low Risk\n"
         "Amber (#F59E0B) - Medium\n"
         "Red (#EF4444) - High Risk"),
        ("Responsive Design", "Mobile-first approach with\n"
         "Tailwind CSS v4. Adapts\nseamlessly to desktop,\ntablet, and mobile views."),
    ]

    for i, (title, desc) in enumerate(design_elements):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.3
        y = 1.8 + row * 2.8
        add_card(slide, x, y, 5.8, 2.4)
        add_text_box(slide, x + 0.3, y + 0.15, 5.2, 0.5, title,
                     font_size=20, bold=True, color=ACCENT_GREEN)
        add_text_box(slide, x + 0.3, y + 0.7, 5.2, 1.6, desc,
                     font_size=14, color=LIGHT_GRAY)


def slide_impact(prs):
    """Slide 16: Impact and Use Cases."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 5, 1.2,
                 "Impact and\nUse Cases", font_size=42, bold=True)
    add_text_box(slide, 0.5, 1.6, 5, 0.5,
                 "Real-world benefits of early patient risk prediction",
                 font_size=14, color=LIGHT_GRAY)

    cards = [
        ("Reduced Critical Events", "Early risk detection helps prevent\n"
         "sudden deterioration, reducing ICU\n"
         "transfers, code-blue events, and\n"
         "emergency interventions."),
        ("Operational Impact", "Better prioritization of high-risk\n"
         "patients reduces length of stay,\n"
         "readmissions, and optimizes staff\n"
         "workload and resources."),
        ("Patient Impact", "Proactive monitoring improves patient\n"
         "outcomes, safety, and confidence\n"
         "through timely clinical attention."),
        ("Use Cases", "Applicable across wards, ICUs,\n"
         "emergency units, and post-operative\n"
         "care for continuous patient\n"
         "risk monitoring."),
    ]

    positions = [(6.5, 0.5), (9.5, 0.5), (6.5, 3.8), (9.5, 3.8)]
    for (title, desc), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 3.2, 3.0)
        add_text_box(slide, x + 0.2, y + 0.2, 2.8, 0.5, title,
                     font_size=16, bold=True)
        add_text_box(slide, x + 0.2, y + 0.8, 2.8, 2.0, desc,
                     font_size=12, color=LIGHT_GRAY)


def slide_future_scope(prs):
    """Slide 17: Future Scope & Scalability."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 0.3, 6, 1.2,
                 "Future Scope &\nScalability", font_size=42, bold=True)
    add_text_box(slide, 7.0, 0.4, 5.5, 0.5,
                 "Strategic initiatives for expanding platform capabilities",
                 font_size=14, color=LIGHT_GRAY)

    items = [
        "Wearable Integration",
        "Predictive Population Health",
        "Clinical Validation",
        "Scalable Deployment",
        "Population Level Insights",
        "Remote Monitoring",
        "Scalable AI Platform",
        "Beyond Hospital Care",
        "Real World Validation",
    ]

    for i, item in enumerate(items):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 2.0 + row * 1.6
        add_card(slide, x, y, 3.8, 1.2)
        add_text_box(slide, x + 0.2, y + 0.3, 3.4, 0.6, item,
                     font_size=16, bold=False, color=WHITE,
                     alignment=PP_ALIGN.CENTER)


def slide_thank_you(prs):
    """Slide 18: Thank You with live link."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text_box(slide, 0.5, 3.0, 12, 2.0,
                 "Thank You", font_size=72, bold=True)

    txBox, tf = add_text_box(slide, 0.5, 5.2, 12, 0.5, "", font_size=16)
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Team: "
    run1.font.bold = False
    run1.font.size = Pt(16)
    run1.font.color.rgb = LIGHT_GRAY
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = "Tech Titans"
    run2.font.bold = True
    run2.font.size = Pt(16)
    run2.font.color.rgb = WHITE
    run2.font.name = "Calibri"

    add_text_box(slide, 0.5, 5.6, 12, 0.4,
                 "Pandit Deendayal Energy University",
                 font_size=16, color=LIGHT_GRAY)

    # Divider line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.2),
        Inches(12), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(100, 100, 100)
    shape.line.fill.background()

    # Live link
    add_text_box(slide, 0.5, 6.3, 12, 0.4,
                 "Thank you for your time and consideration.",
                 font_size=14, color=LIGHT_GRAY)

    # Live hosted link - PROMINENT
    add_card(slide, 3.0, 6.7, 7.5, 0.6)
    add_text_box(slide, 3.2, 6.72, 7.1, 0.5,
                 "Live Website:  https://cortex-health-app.vercel.app",
                 font_size=16, bold=True, color=ACCENT_GREEN,
                 alignment=PP_ALIGN.CENTER)


# ─────────────────────────── MAIN ───────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all slides
    slide_cover(prs)                                              # 1
    slide_section_title(prs,                                      # 2
        "Predicting Patient Health Risk\nThrough Proactive Monitoring",
        "Phase 2: Fully working website - deployed and live")
    slide_problem_recap(prs)                                      # 3
    slide_solution_overview(prs)                                  # 4
    slide_tech_stack(prs)                                         # 5
    slide_system_architecture(prs)                                # 6
    slide_ml_model(prs)                                           # 7
    slide_safety_overrides(prs)                                   # 8
    slide_dashboard_features(prs)                                 # 9
    slide_auth_security(prs)                                      # 10
    slide_api_endpoints(prs)                                      # 11
    slide_deployment(prs)                                         # 12
    slide_data_flow(prs)                                          # 13
    slide_model_performance(prs)                                  # 14
    slide_ui_design(prs)                                          # 15
    slide_impact(prs)                                             # 16
    slide_future_scope(prs)                                       # 17
    slide_thank_you(prs)                                          # 18

    prs.save(OUTPUT_PATH)
    print(f"Phase 2 PPT generated: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
